#!/usr/bin/env python3
"""将 Markdown 格式的 PPT 内容转换为 PPTX 文件 - 优化版"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def parse_slide(slide_text):
    """解析单页幻灯片内容"""
    lines = [l for l in slide_text.split('\n') if l.strip()]
    if not lines:
        return None, []
    
    title = ""
    content_items = []
    
    # 解析标题
    for line in lines:
        if line.startswith('# '):
            title = line[2:].strip()
            break
        elif line.startswith('## '):
            title = line[3:].strip()
            break
    
    # 解析内容
    in_code = False
    code_lines = []
    table_rows = []
    in_table = False
    
    for line in lines:
        line_stripped = line.strip()
        
        # 跳过标题行
        if line_stripped.startswith('#'):
            continue
        
        # 代码块处理
        if line_stripped.startswith('```'):
            if in_code:
                content_items.append(('code', '\n'.join(code_lines)))
                code_lines = []
                in_code = False
            else:
                in_code = True
            continue
        
        if in_code:
            code_lines.append(line)
            continue
        
        # 表格处理
        if '|' in line_stripped and line_stripped.startswith('|') and line_stripped.endswith('|'):
            cells = [c.strip() for c in line.split('|')[1:-1]]
            if len(cells) > 0 and not all(c.startswith('-') or c.startswith(':') for c in cells):
                if not in_table:
                    in_table = True
                    table_rows = []
                table_rows.append(cells)
                continue
        
        if in_table and table_rows:
            content_items.append(('table', table_rows))
            table_rows = []
            in_table = False
        
        # 列表项
        if line_stripped.startswith('- ') or line_stripped.startswith('* '):
            content_items.append(('bullet', line_stripped[2:].strip()))
        # 普通文本
        elif line_stripped and not line_stripped.startswith('|'):
            content_items.append(('text', line_stripped))
    
    if table_rows:
        content_items.append(('table', table_rows))
    
    return title, content_items

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    with open('/Users/hirol/.openclaw/workspace/ai-sharing-ppt.md', 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides_content = content.split('---')
    slide_num = 0
    empty_slides = 0
    
    for i, slide_text in enumerate(slides_content):
        slide_text = slide_text.strip()
        if not slide_text:
            continue
        
        title, content_items = parse_slide(slide_text)
        
        # 跳过无标题且无内容的页面
        if not title and not content_items:
            empty_slides += 1
            continue
        
        slide_num += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 添加标题
        y_position = 0.4
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(12.333), Inches(0.8))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            y_position = 1.3
        
        # 添加内容
        for item_type, content in content_items:
            if y_position > 6.8:
                break
            
            if item_type == 'table' and content:
                rows = len(content)
                cols = max(len(row) for row in content) if content else 1
                if cols == 0:
                    cols = 1
                
                table_height = min(rows * 0.5, 3.5)
                table = slide.shapes.add_table(
                    rows, cols, Inches(0.5), Inches(y_position),
                    Inches(12.333), Inches(table_height)
                ).table
                
                for row_idx, row in enumerate(content):
                    for col_idx, cell_text in enumerate(row):
                        if col_idx < cols:
                            cell = table.cell(row_idx, col_idx)
                            cell.text = cell_text
                            for para in cell.text_frame.paragraphs:
                                para.font.size = Pt(13)
                
                y_position += table_height + 0.2
            
            elif item_type == 'code':
                if content and len(content) < 1200:
                    code_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(y_position),
                        Inches(12.333), Inches(1.8)
                    )
                    code_frame = code_box.text_frame
                    code_para = code_frame.paragraphs[0]
                    code_para.text = content
                    code_para.font.size = Pt(10)
                    code_para.font.name = 'Courier New'
                    y_position += 2.0
            
            elif item_type == 'bullet':
                if content:
                    box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(y_position),
                        Inches(12.333), Inches(0.35)
                    )
                    frame = box.text_frame
                    para = frame.paragraphs[0]
                    para.text = f"• {content}"
                    para.font.size = Pt(18)
                    y_position += 0.4
            
            elif item_type == 'text':
                if content and len(content) < 150:
                    box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(y_position),
                        Inches(12.333), Inches(0.35)
                    )
                    frame = box.text_frame
                    para = frame.paragraphs[0]
                    para.text = content
                    para.font.size = Pt(18)
                    y_position += 0.4
    
    output_path = '/Users/hirol/.openclaw/workspace/AI 赋能产研效率提升分享.pptx'
    prs.save(output_path)
    
    print(f"\n✅ PPT 已生成：{output_path}")
    print(f"📊 共 {slide_num} 页")
    if empty_slides > 0:
        print(f"⚠️  跳过 {empty_slides} 页空白内容")
    
    return output_path

if __name__ == '__main__':
    create_ppt()
